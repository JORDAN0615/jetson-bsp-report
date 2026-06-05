#!/usr/bin/env python3
"""Create MIC-741 verification PPTX pages from a white logo template."""

from __future__ import annotations

import argparse
import re
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt


DEFAULT_TEMPLATE = Path("assets/blank_logo_template.pptx")

TITLE_LEFT = 457200
TITLE_TOP = 87313
TITLE_WIDTH = 8229600
TITLE_HEIGHT = 550000

BODY_LEFT = 457200
BODY_TOP = 760000
BODY_WIDTH = 8229600
BODY_HEIGHT = 2200000

PAGE_LEFT = 9
PAGE_TOP = 4749901
PAGE_WIDTH = 548700
PAGE_HEIGHT = 393600

PACKET_IMAGE_LEFT = 1260546
PACKET_IMAGE_TOP = 3005835
PACKET_IMAGE_WIDTH = 6622908

BLOCK_LEFT = 100000
BLOCK_TOP = 2300000
BLOCK_WIDTH = 7400000
BLOCK_ROW_STEP = 395000

TEXT_COLOR = RGBColor(0, 0, 0)
TITLE_COLOR = RGBColor(17, 50, 99)

PACKET_TITLE = "[MIC-741] USB3C1, USB3C2, USB 3.0 Ports (M8, M9)"
BLOCK_TITLE = "[MIC-741] USB3C1, USB3C2, USB 3.0 (M8, M9)"

PACKET_BODY = """Please open a terminal and run below commands. If USB 3.0 dongle is connected successfully, the output should be 9.

sudo cat /sys/bus/usb/devices/2-3.1/bMaxPacketSize0 | grep 9
sudo cat /sys/bus/usb/devices/2-3.2/bMaxPacketSize0 | grep 9
sudo cat /sys/bus/usb/devices/2-3.3/bMaxPacketSize0 | grep 9
sudo cat /sys/bus/usb/devices/2-3.4/bMaxPacketSize0 | grep 9
sudo cat /sys/bus/usb/devices/2-1.4/bMaxPacketSize0 | grep 9
sudo cat /sys/bus/usb/devices/2-1.3/bMaxPacketSize0 | grep 9"""

BLOCK_BODY = """Please open a terminal and run below commands. If USB 3.0 dongle is connected successfully, a block/sdX device path should be shown.

sudo find /sys/bus/usb/devices/2-3.1/ | grep block/sd.$
sudo find /sys/bus/usb/devices/2-3.2/ | grep block/sd.$
sudo find /sys/bus/usb/devices/2-3.3/ | grep block/sd.$
sudo find /sys/bus/usb/devices/2-3.4/ | grep block/sd.$
sudo find /sys/bus/usb/devices/2-1.3/ | grep block/sd.$
sudo find /sys/bus/usb/devices/2-1.4/ | grep block/sd.$"""


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide_id)


def keep_first_slides(prs: Presentation, count: int) -> None:
    if len(prs.slides) < count:
        raise ValueError(f"template needs at least {count} slide(s), got {len(prs.slides)}")
    for index in reversed(range(len(prs.slides))):
        if index >= count:
            delete_slide(prs, index)


def add_text(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    size: int,
    bold: bool = False,
    color: RGBColor = TEXT_COLOR,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    lines = text.splitlines() or [""]
    first = frame.paragraphs[0]
    first.text = lines[0]
    for line in lines[1:]:
        paragraph = frame.add_paragraph()
        paragraph.text = line

    for paragraph in frame.paragraphs:
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color


def add_page_number(slide, number: int) -> None:
    add_text(slide, PAGE_LEFT, PAGE_TOP, PAGE_WIDTH, PAGE_HEIGHT, str(number), 18, bold=True, color=TITLE_COLOR)


def add_packet_page(slide, screenshot: Path) -> None:
    add_text(slide, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT, PACKET_TITLE, 20, bold=True, color=TITLE_COLOR)
    add_text(slide, BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_HEIGHT, PACKET_BODY, 12)
    slide.shapes.add_picture(str(screenshot), PACKET_IMAGE_LEFT, PACKET_IMAGE_TOP, width=PACKET_IMAGE_WIDTH)
    add_page_number(slide, 22)


def add_block_page(slide, screenshots: list[Path]) -> None:
    if len(screenshots) != 6:
        raise ValueError(f"USB block page requires exactly 6 screenshots, got {len(screenshots)}")
    add_text(slide, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT, BLOCK_TITLE, 20, bold=True, color=TITLE_COLOR)
    add_text(slide, BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_HEIGHT, BLOCK_BODY, 12)
    positions = [
        (BLOCK_LEFT, BLOCK_TOP),
        (BLOCK_LEFT, BLOCK_TOP + BLOCK_ROW_STEP),
        (BLOCK_LEFT, BLOCK_TOP + BLOCK_ROW_STEP * 2),
        (BLOCK_LEFT, BLOCK_TOP + BLOCK_ROW_STEP * 3),
        (BLOCK_LEFT, BLOCK_TOP + BLOCK_ROW_STEP * 4),
        (BLOCK_LEFT, BLOCK_TOP + BLOCK_ROW_STEP * 5),
    ]
    for image_path, (x, y) in zip(screenshots, positions):
        slide.shapes.add_picture(str(image_path), x, y, width=BLOCK_WIDTH)
    add_page_number(slide, 23)


def create_report(template: Path, screenshot: Path, output: Path) -> None:
    prs = Presentation(str(template))
    keep_first_slides(prs, 1)
    add_packet_page(prs.slides[0], screenshot)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def create_two_page_report(template: Path, packet_screenshot: Path, block_screenshots: list[Path], output: Path) -> None:
    prs = Presentation(str(template))
    keep_first_slides(prs, 2)
    add_packet_page(prs.slides[0], packet_screenshot)
    add_block_page(prs.slides[1], block_screenshots)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def count_slides(pptx_path: Path) -> int:
    with zipfile.ZipFile(pptx_path) as archive:
        return len([n for n in archive.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)])


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--template", default=str(DEFAULT_TEMPLATE), help="White logo template PPTX path.")
    parser.add_argument("--screenshot", required=True, help="Generated terminal PNG path.")
    parser.add_argument(
        "--block-screenshots",
        nargs="*",
        default=[],
        help="Six generated terminal PNG paths for slide 23 USB block checks.",
    )
    parser.add_argument("--output", required=True, help="Output PPTX path.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = Path(args.output)
    if args.block_screenshots:
        create_two_page_report(
            Path(args.template),
            Path(args.screenshot),
            [Path(path) for path in args.block_screenshots],
            output,
        )
        expected_slides = 2
    else:
        create_report(Path(args.template), Path(args.screenshot), output)
        expected_slides = 1
    slides = count_slides(output)
    if slides != expected_slides:
        raise SystemExit(f"Expected {expected_slides} slide(s), got {slides}.")
    print(f"Wrote {output} with {slides} slide.")


if __name__ == "__main__":
    main()
