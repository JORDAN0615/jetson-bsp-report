#!/usr/bin/env python3
"""Create MIC-741 verification PPTX pages by replacing template evidence."""

from __future__ import annotations

import argparse
import re
import zipfile
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.shapes.picture import Picture


RED_RGB = "FF0000"


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide_id)


def keep_only_slide(prs: Presentation, one_based_slide: int) -> None:
    keep = one_based_slide - 1
    if keep < 0 or keep >= len(prs.slides):
        raise ValueError(f"source slide {one_based_slide} is outside deck with {len(prs.slides)} slides")
    for index in reversed(range(len(prs.slides))):
        if index != keep:
            delete_slide(prs, index)


def keep_only_slides(prs: Presentation, one_based_slides: list[int]) -> None:
    keep = {slide_number - 1 for slide_number in one_based_slides}
    if any(index < 0 or index >= len(prs.slides) for index in keep):
        raise ValueError(f"source slides {one_based_slides} are outside deck with {len(prs.slides)} slides")
    for index in reversed(range(len(prs.slides))):
        if index not in keep:
            delete_slide(prs, index)


def iter_pictures(slide) -> Iterable[Picture]:
    for shape in slide.shapes:
        if isinstance(shape, Picture):
            yield shape


def find_widest_picture(slide) -> Picture:
    pictures = list(iter_pictures(slide))
    if not pictures:
        raise ValueError("No picture shapes found on source slide.")
    return max(pictures, key=lambda shape: shape.width * shape.height)


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def replace_picture_image(picture: Picture, image_path: Path) -> None:
    image_part, rel_id = picture.part.get_or_add_image_part(str(image_path))
    blip = picture._pic.blipFill.blip
    embed_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
    blip.set(embed_attr, rel_id)


def overlaps(a, b) -> bool:
    return not (
        a.left + a.width <= b.left
        or b.left + b.width <= a.left
        or a.top + a.height <= b.top
        or b.top + b.height <= a.top
    )


def remove_template_evidence_overlays(slide, picture: Picture) -> None:
    """Remove the template's old red result box when the new PNG already has it."""
    for shape in list(slide.shapes):
        if shape is picture or not overlaps(shape, picture):
            continue
        if hasattr(shape, "text") and shape.text:
            continue
        try:
            line_rgb = str(shape.line.color.rgb)
        except (AttributeError, TypeError):
            line_rgb = ""
        is_narrow_result_box = shape.width < picture.width * 0.12 and shape.height > picture.height * 0.5
        if line_rgb == RED_RGB and is_narrow_result_box:
            remove_shape(shape)


def create_report(template: Path, screenshot: Path, output: Path, source_slide: int) -> None:
    prs = Presentation(str(template))
    keep_only_slide(prs, source_slide)
    slide = prs.slides[0]
    picture = find_widest_picture(slide)
    remove_template_evidence_overlays(slide, picture)
    replace_picture_image(picture, screenshot)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def add_usb_block_screenshots(slide, screenshots: list[Path]) -> None:
    if len(screenshots) != 6:
        raise ValueError(f"USB block page requires exactly 6 screenshots, got {len(screenshots)}")
    original_picture = find_widest_picture(slide)
    remove_shape(original_picture)

    left_x = 323257
    col_gap = 190000
    top_y = 3050000
    width = 4060000
    row_step = 535000
    positions = [
        (left_x, top_y),
        (left_x + width + col_gap, top_y),
        (left_x, top_y + row_step),
        (left_x + width + col_gap, top_y + row_step),
        (left_x, top_y + row_step * 2),
        (left_x + width + col_gap, top_y + row_step * 2),
    ]
    for image_path, (x, y) in zip(screenshots, positions):
        slide.shapes.add_picture(str(image_path), x, y, width=width)


def create_two_page_report(
    template: Path,
    packet_screenshot: Path,
    block_screenshots: list[Path],
    output: Path,
    packet_slide: int,
    block_slide: int,
) -> None:
    prs = Presentation(str(template))
    keep_only_slides(prs, [packet_slide, block_slide])

    packet_page = prs.slides[0]
    packet_picture = find_widest_picture(packet_page)
    remove_template_evidence_overlays(packet_page, packet_picture)
    replace_picture_image(packet_picture, packet_screenshot)

    block_page = prs.slides[1]
    add_usb_block_screenshots(block_page, block_screenshots)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def count_slides(pptx_path: Path) -> int:
    with zipfile.ZipFile(pptx_path) as archive:
        return len([n for n in archive.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)])


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--template", required=True, help="Template PPTX path.")
    parser.add_argument("--screenshot", required=True, help="Generated terminal PNG path.")
    parser.add_argument(
        "--block-screenshots",
        nargs="*",
        default=[],
        help="Six generated terminal PNG paths for slide 23 USB block checks.",
    )
    parser.add_argument("--output", required=True, help="Output one-slide PPTX path.")
    parser.add_argument("--source-slide", type=int, default=22, help="One-based source slide number.")
    parser.add_argument("--block-source-slide", type=int, default=23, help="One-based USB block source slide number.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = Path(args.output)
    if args.block_screenshots:
        create_two_page_report(
            Path(args.template),
            Path(args.screenshot),
            [Path(path) for path in args.block_screenshots],
            output,
            args.source_slide,
            args.block_source_slide,
        )
        expected_slides = 2
    else:
        create_report(Path(args.template), Path(args.screenshot), output, args.source_slide)
        expected_slides = 1
    slides = count_slides(output)
    if slides != expected_slides:
        raise SystemExit(f"Expected {expected_slides} slide(s), got {slides}.")
    print(f"Wrote {output} with {slides} slide.")


if __name__ == "__main__":
    main()
